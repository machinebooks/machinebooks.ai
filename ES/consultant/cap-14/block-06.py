# Extraído de: LibroConsultor/cap-14-reporting.md
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def generar_presentacion_ejecutiva(
    proyecto: ProyectoReporting,
    resumen_ejecutivo: str,
    hallazgos_top: list[Hallazgo],
    recomendaciones_priorizadas: str,
    plantilla_pptx: str,
    output_path: str
) -> str:
    """Genera presentación ejecutiva de 12-15 diapositivas."""

    prs = Presentation(plantilla_pptx)

    # Diapositiva 1: Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = proyecto.nombre_proyecto
    slide.placeholders[1].text = (
        f"{proyecto.cliente} | {proyecto.fecha_fin}"
    )

    # Diapositiva 2: Resumen ejecutivo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resumen ejecutivo"
    # Truncar a los primeros 3 párrafos para la diapositiva
    parrafos = resumen_ejecutivo.split("\n\n")[:3]
    body = slide.placeholders[1]
    body.text = "\n\n".join(parrafos)

    # Diapositiva 3: Distribución de hallazgos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hallazgos por severidad"
    conteo = {}
    for h in proyecto.hallazgos:
        conteo[h.severidad.value] = conteo.get(h.severidad.value, 0) + 1
    body = slide.placeholders[1]
    body.text = "\n".join(
        f"{sev.upper()}: {n} hallazgos" for sev, n in conteo.items()
    )

    # Diapositivas 4-8: Top hallazgos (uno por diapositiva)
    for h in hallazgos_top[:5]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = h.titulo
        body = slide.placeholders[1]
        body.text = (
            f"Severidad: {h.severidad.value.upper()}\n\n"
            f"{h.descripcion}\n\n"
            f"Impacto: {h.impacto_negocio}\n\n"
            f"Recomendación: {h.recomendacion}"
        )

    # Diapositiva 9: Roadmap resumido
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roadmap de implementación"
    body = slide.placeholders[1]
    body.text = recomendaciones_priorizadas[:1500]  # Resumen

    # Diapositiva final: Próximos pasos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Próximos pasos"
    body = slide.placeholders[1]
    body.text = (
        "1. Validar hallazgos con equipos técnicos (semana 1)\n"
        "2. Priorizar recomendaciones con presupuesto (semana 2)\n"
        "3. Asignar responsables y plazos (semana 3)\n"
        "4. Sesión de arranque del plan de acción (semana 4)"
    )

    prs.save(output_path)
    return output_path
